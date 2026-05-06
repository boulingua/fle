"""
Generate per-unit Materials placeholders (presentation .pptx, worksheet .pdf,
thumbnail .png × 2) and inject `presentation:` / `worksheet:` frontmatter into
every unit .md under content/track_*/units/.

Idempotent: regenerating only overwrites placeholders + updates frontmatter
keys; does not touch article body. Safe to re-run in CI.
"""
from __future__ import annotations

import io
import pathlib
import re
import sys
from typing import Optional

REPO = pathlib.Path(__file__).resolve().parent.parent
CONTENT = REPO / "content"
STATIC_BASE = REPO / "static" / "materials"
PRES_DIR = STATIC_BASE / "presentations"
WORK_DIR = STATIC_BASE / "worksheets"

YAML_RE = re.compile(r"^---\n(.*?)\n---\n", re.DOTALL)


def list_unit_pages() -> list[pathlib.Path]:
    return sorted(p for p in CONTENT.rglob("units/*.md") if not p.name.startswith("_"))


def parse_fm(text: str) -> tuple[dict, str, str]:
    m = YAML_RE.match(text)
    if not m:
        return {}, "", text
    fm_text = m.group(1)
    fm: dict = {}
    for line in fm_text.splitlines():
        sm = re.match(r'^([A-Za-z][\w-]*):\s*(.*)$', line)
        if sm:
            key, val = sm.group(1), sm.group(2).strip()
            if val and not val.startswith("|") and not val.startswith(">"):
                fm[key] = val.strip().strip('"')
    return fm, fm_text, text[m.end():]


def upsert_fm_block(fm_text: str, key: str, lines: list[str]) -> str:
    """Replace or insert a block-style key (`key:\n  …`) into fm_text."""
    pattern = re.compile(
        rf"^{re.escape(key)}:\s*\n(?:[ \t]+\S.*\n?)*",
        re.MULTILINE,
    )
    block = key + ":\n" + "\n".join("  " + l for l in lines) + "\n"
    if pattern.search(fm_text):
        return pattern.sub(block, fm_text, count=1)
    if not fm_text.endswith("\n"):
        fm_text += "\n"
    return fm_text + block


def make_pptx(out: pathlib.Path, title: str, subtitle: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    layout = prs.slide_layouts[5]  # title only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(3))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.text = "Placeholder — à remplacer par la version finale."
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(18)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))


def make_pdf(out: pathlib.Path, title: str, subtitle: str) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    out.parent.mkdir(parents=True, exist_ok=True)
    doc = SimpleDocTemplate(
        str(out), pagesize=A4,
        leftMargin=2.5 * cm, rightMargin=2.5 * cm,
        topMargin=2.5 * cm, bottomMargin=2.5 * cm,
        author="S. Le Boulanger", title=title,
    )
    styles = getSampleStyleSheet()
    story = [
        Paragraph(title, styles["Title"]),
        Spacer(1, 0.5 * cm),
        Paragraph(subtitle, styles["Heading3"]),
        Spacer(1, 1 * cm),
        Paragraph(
            "Fiche d'exercices — placeholder. Sera remplacée par la version finale.",
            styles["BodyText"],
        ),
    ]
    doc.build(story)


def render_pdf_first_page(pdf: pathlib.Path, png: pathlib.Path,
                          width: int = 480) -> bool:
    """Render page 1 of `pdf` to `png`. Pure-Python via pypdfium2 — no
    poppler / cairo system deps. Returns True on success."""
    try:
        import pypdfium2 as pdfium
    except ImportError:
        return False
    try:
        doc = pdfium.PdfDocument(str(pdf))
        if len(doc) == 0:
            return False
        page = doc[0]
        # Match width to thumb width; scale = width / page.width_in_points
        scale = width / page.get_width()
        bitmap = page.render(scale=scale)
        img = bitmap.to_pil()
        png.parent.mkdir(parents=True, exist_ok=True)
        img.save(png, "PNG")
        return True
    except Exception as exc:  # pragma: no cover
        print(f"  pdfium failed for {pdf.name}: {exc}", file=sys.stderr)
        return False


def _find_libreoffice() -> Optional[str]:
    import shutil
    for cmd in ("soffice", "libreoffice"):
        path = shutil.which(cmd)
        if path:
            return path
    # Common Windows install paths
    for p in (
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ):
        if pathlib.Path(p).exists():
            return p
    return None


def render_pptx_first_slide(pptx: pathlib.Path, png: pathlib.Path,
                            width: int = 480) -> bool:
    """Convert `pptx` → PDF via LibreOffice, then render page 1 to PNG.
    Returns False (caller falls back to PIL placeholder) if LibreOffice
    isn't installed or conversion fails."""
    soffice = _find_libreoffice()
    if not soffice:
        return False
    import subprocess
    import tempfile
    with tempfile.TemporaryDirectory() as td:
        td_path = pathlib.Path(td)
        try:
            subprocess.run(
                [soffice, "--headless", "--norestore",
                 "--convert-to", "pdf",
                 "--outdir", str(td_path), str(pptx)],
                check=True, capture_output=True, timeout=60,
            )
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired) as exc:
            print(f"  soffice failed for {pptx.name}: {exc}", file=sys.stderr)
            return False
        candidates = list(td_path.glob(f"{pptx.stem}.pdf"))
        if not candidates:
            return False
        return render_pdf_first_page(candidates[0], png, width=width)


def make_thumbnail(out: pathlib.Path, title: str, kind: str, accent: str) -> None:
    from PIL import Image, ImageDraw, ImageFont

    W, H = 480, 320
    img = Image.new("RGB", (W, H), "white")
    draw = ImageDraw.Draw(img)
    # Coloured stripe top
    draw.rectangle([0, 0, W, 60], fill=accent)
    try:
        font_big = ImageFont.truetype("DejaVuSans-Bold.ttf", 22)
        font_small = ImageFont.truetype("DejaVuSans.ttf", 14)
    except OSError:
        font_big = ImageFont.load_default()
        font_small = ImageFont.load_default()
    draw.text((20, 18), kind, fill="white", font=font_big)
    # Title wrapped
    text = title
    lines: list[str] = []
    cur = ""
    for word in text.split():
        candidate = (cur + " " + word).strip()
        if draw.textlength(candidate, font=font_small) > W - 40 and cur:
            lines.append(cur)
            cur = word
        else:
            cur = candidate
    if cur:
        lines.append(cur)
    y = 100
    for line in lines[:8]:
        draw.text((20, y), line, fill="black", font=font_small)
        y += 22
    draw.text((20, H - 28), "Placeholder", fill="#888", font=font_small)
    draw.rectangle([0, 0, W - 1, H - 1], outline="#ddd")
    out.parent.mkdir(parents=True, exist_ok=True)
    img.save(out, "PNG")


def derive_tags(fm: dict, path: pathlib.Path) -> list[str]:
    tags: list[str] = []
    track = fm.get("track")
    if track:
        tags.append(f"filiere-{track.lower()}")
    klasse = fm.get("klassenstufe")
    if klasse:
        tags.append(f"classe-{klasse}")
    niveau = fm.get("niveau")
    if niveau:
        tags.append(f"niveau-{niveau.lower()}")
    if not tags:
        # Fall back to path-derived
        for part in path.parts:
            if part.startswith("track_"):
                tags.append(part)
                break
    return tags


def derive_short_label(fm: dict) -> str:
    klasse = fm.get("klassenstufe", "")
    niveau = fm.get("niveau", "")
    parts = []
    if klasse:
        parts.append(f"Cl. {klasse}")
    if niveau:
        parts.append(f"Niveau {niveau}")
    return " · ".join(parts)


def process_unit(md_path: pathlib.Path) -> dict:
    text = md_path.read_text(encoding="utf-8")
    fm, fm_text, body = parse_fm(text)
    slug = md_path.stem  # filename stem, stable
    title = fm.get("title", slug)
    subtitle = fm.get("subtitle", derive_short_label(fm) or "FLE — placeholder")
    rel = md_path.relative_to(REPO).as_posix()

    pres_pptx = PRES_DIR / f"{slug}.pptx"
    pres_png = PRES_DIR / f"{slug}.png"
    work_pdf = WORK_DIR / f"{slug}.pdf"
    work_png = WORK_DIR / f"{slug}.png"

    make_pptx(pres_pptx, title, subtitle)
    make_pdf(work_pdf, title, subtitle)
    # Real thumbnails first; fall back to PIL placeholder if the toolchain
    # is missing (LibreOffice for PPTX, pypdfium2 for PDF).
    if not render_pptx_first_slide(pres_pptx, pres_png):
        make_thumbnail(pres_png, title, "Présentation", "#1a73e8")
    if not render_pdf_first_page(work_pdf, work_png):
        make_thumbnail(work_png, title, "Fiche d'exercices", "#2e7d32")

    tags = derive_tags(fm, md_path.relative_to(REPO))
    # Hardcode the production URL prefix. Hugo's relURL/absURL do NOT
    # prepend the baseURL path component to leading-slash inputs, so a
    # frontmatter value of "/materials/foo" renders as href="/materials/
    # foo" — site-root resolved, missing the /fle/ prefix → 404. Baking
    # /fle/ into the data here keeps both server templates and the
    # client-side JS (which reads graph.json verbatim) honest. If the
    # site ever moves off /fle/, regenerate.
    PREFIX = "/fle"
    pres_url = f"{PREFIX}/materials/presentations/{slug}.pptx"
    pres_thumb = f"{PREFIX}/materials/presentations/{slug}.png"
    work_url = f"{PREFIX}/materials/worksheets/{slug}.pdf"
    work_thumb = f"{PREFIX}/materials/worksheets/{slug}.png"

    new_fm = upsert_fm_block(fm_text, "presentation", [
        f'file: "{pres_url}"',
        f'thumbnail: "{pres_thumb}"',
    ])
    new_fm = upsert_fm_block(new_fm, "worksheet", [
        f'file: "{work_url}"',
        f'thumbnail: "{work_thumb}"',
    ])
    if "tags:" not in new_fm and tags:
        new_fm += "tags:\n" + "\n".join(f"  - {t}" for t in tags) + "\n"

    new_text = "---\n" + new_fm.rstrip("\n") + "\n---\n" + body
    md_path.write_text(new_text, encoding="utf-8")
    return {"slug": slug, "title": title, "tags": tags}


def main() -> int:
    units = list_unit_pages()
    if not units:
        print("no unit pages found", file=sys.stderr)
        return 1
    PRES_DIR.mkdir(parents=True, exist_ok=True)
    WORK_DIR.mkdir(parents=True, exist_ok=True)
    for i, p in enumerate(units, 1):
        info = process_unit(p)
        if i % 25 == 0 or i == len(units):
            print(f"  [{i}/{len(units)}] {info['slug']}")
    print(f"Done — {len(units)} units; "
          f"{len(list(PRES_DIR.glob('*.pptx')))} pptx, "
          f"{len(list(WORK_DIR.glob('*.pdf')))} pdf, "
          f"{len(list(PRES_DIR.glob('*.png')))} pres-thumbs, "
          f"{len(list(WORK_DIR.glob('*.png')))} work-thumbs.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
